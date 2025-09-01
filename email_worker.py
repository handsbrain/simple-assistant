# email_worker.py
from __future__ import annotations
import os, sys, time, json, signal, threading, re, html, random, base64
from pathlib import Path
from typing import Dict, Any, List, Optional
from datetime import datetime, timezone

# -------- Vector memory ----------
from core.vector_mem import add_memory, search_memory, all_stats
try:
    from core.vector_mem import all_stats as mem_all_stats
    _HAS_MEM = True
except Exception:
    mem_all_stats = None
    _HAS_MEM = False

# -------- LLM ----------
from core.llm import ask_chatgpt

# -------- Config ----------
POLL_SECONDS      = int(os.getenv("POLL_SECONDS", "20"))
DRY_RUN           = os.getenv("DRY_RUN", "0").lower() in ("1", "true", "yes", "y")
HEALTH_PORT       = int(os.getenv("HEALTH_PORT", "8080"))
MAX_PROMPT_CHARS  = int(os.getenv("MAX_PROMPT_CHARS", "20000"))
SIGNATURE_ENV     = os.getenv("ASSISTANT_SIGNATURE", "")          # optional
SIGNATURE_FILE    = os.getenv("ASSISTANT_SIGNATURE_FILE", "")     # optional (path)

# Allowlist (email or bare domain like example.com)
_raw_allow = [x.strip().lower() for x in os.getenv("ALLOWLIST_SENDERS", "").split(",") if x.strip()]
ALLOWLIST_EMAILS  = {a for a in _raw_allow if "@" in a}
ALLOWLIST_DOMAINS = {a.split("@",1)[1] for a in _raw_allow if "@" in a} | {a for a in _raw_allow if "@" not in a}

STATE_DIR   = Path(os.getenv("STATE_DIR", "state"))
STATE_DIR.mkdir(parents=True, exist_ok=True)
STATE_PATH  = STATE_DIR / "processed.json"

# -------- Attachments config ----------
ATTACH_ENABLE     = os.getenv("ATTACH_ENABLE", "1").lower() in ("1","true","yes","y")
ATTACH_MAX_COUNT  = int(os.getenv("ATTACH_MAX_COUNT", "50"))
ATTACH_MAX_MB     = int(os.getenv("ATTACH_MAX_MB", "30"))
ATTACH_EXTS       = {e.strip().lower() for e in (os.getenv("ATTACH_EXTS", "pdf,docx,pptx,xlsx,txt,csv,png,jpg,jpeg,tiff,bmp,webp").split(",")) if e.strip()}
ATTACH_SUMMARY_MAX_CHARS = int(os.getenv("ATTACH_SUMMARY_MAX_CHARS", "2000"))

# OCR (Tesseract) options
ATTACH_OCR        = os.getenv("ATTACH_OCR", "1").lower() in ("1","true","yes","y")
OCR_PAGES_MAX     = int(os.getenv("ATTACH_OCR_PAGES_MAX", "10"))
OCR_LANG          = os.getenv("ATTACH_OCR_LANG", "eng").strip() or "eng"
OCR_DPI           = int(os.getenv("ATTACH_OCR_DPI", "200"))

# -------- Utilities ----------
def html_to_text(s: str) -> str:
    s = re.sub(r"(?is)<(script|style).*?>.*?</\1>", "", s or "")
    s = re.sub(r"(?is)<br\s*/?>", "\n", s)
    s = re.sub(r"(?is)</p>", "\n\n", s)
    s = re.sub(r"(?is)<.*?>", "", s)
    return html.unescape(s)

def validate_email_content(content: str, max_length: int = 100000) -> str:
    """Validate and sanitize email content."""
    if not content:
        return ""
    
    # Check length
    if len(content) > max_length:
        log(f"[WARNING] Email content truncated from {len(content)} to {max_length} characters")
        content = content[:max_length]
    
    # Remove potential prompt injection attempts
    content = re.sub(r'(?i)(ignore|forget|disregard)\s+(previous|all|above)\s+(instructions?|prompts?|rules?)', '', content)
    content = re.sub(r'(?i)(you\s+are\s+now|act\s+as|pretend\s+to\s+be)', '', content)
    
    return content.strip()

def validate_sender(sender_info: Dict[str, Any]) -> bool:
    """Validate sender information."""
    if not sender_info:
        return False
    
    email_address = sender_info.get("emailAddress", {}).get("address", "")
    if not email_address or "@" not in email_address:
        return False
    
    # Basic email format validation
    if not re.match(r'^[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}$', email_address):
        return False
    
    return True

def _strip_subject_lines(s: str) -> str:
    # Remove any leading "Subject:" line the model might generate
    lines = (s or "").splitlines()
    out = []
    for i, L in enumerate(lines):
        if i == 0 and re.match(r"(?i)^\s*subject\s*:", L):
            continue
        out.append(L)
    return "\n".join(out).strip()

def _strip_llm_signature(s: str) -> str:
    # Heuristic: drop a trailing short sign-off like "Best," on a single line
    s = s.rstrip()
    tail = s.splitlines()[-1].strip().lower() if s else ""
    if tail in {"best", "best,", "thanks", "thanks,", "regards", "regards,"}:
        return "\n".join(s.splitlines()[:-1]).rstrip()
    return s

def _text_to_html(s: str) -> str:
    # Normalize CRLF → LF, escape, then convert LF to <br>
    s = (s or "").replace("\r\n", "\n").replace("\r", "\n")
    esc = html.escape(s, quote=False)
    return esc.replace("\n", "<br>")

def _file_ext(name: str, ctype: str) -> str:
    n = (name or "").lower().strip()
    if "." in n:
        return n.rsplit(".", 1)[-1]
    if "/" in (ctype or ""):
        return ctype.split("/")[-1].lower()
    return ""

def _safe_decode_text(data: bytes) -> str:
    try:
        return data.decode("utf-8")
    except Exception:
        return data.decode("latin-1", errors="ignore")

def _extract_text_from_attachment(name: str, content_type: str, data: bytes) -> str:
    if not data:
        return ""
    if len(data) > ATTACH_MAX_MB * 1024 * 1024:
        return ""
    ext = _file_ext(name, content_type)
    try:
        if ext == "pdf":
            import io
            from pypdf import PdfReader
            r = PdfReader(io.BytesIO(data))
            pdf_text = "\n".join((p.extract_text() or "") for p in r.pages)
            if pdf_text.strip():
                return pdf_text
            if ATTACH_OCR:
                try:
                    import io
                    import pypdfium2 as pdfium
                    from PIL import Image
                    import pytesseract
                    pdf = pdfium.PdfDocument(io.BytesIO(data))
                    n = min(len(pdf), OCR_PAGES_MAX)
                    out_lines = []
                    for i in range(n):
                        page = pdf[i]
                        pil_img = page.render(scale=OCR_DPI / 72.0).to_pil()
                        out_lines.append(pytesseract.image_to_string(pil_img, lang=OCR_LANG))
                    return "\n".join(out_lines)
                except Exception:
                    return ""
        if ext == "docx":
            import io
            from docx import Document
            d = Document(io.BytesIO(data))
            return "\n".join(p.text for p in d.paragraphs if p.text)
        if ext == "pptx":
            import io
            from pptx import Presentation
            prs = Presentation(io.BytesIO(data))
            out = []
            for slide in prs.slides:
                for shp in slide.shapes:
                    if hasattr(shp, "text") and shp.text:
                        out.append(shp.text)
            return "\n".join(out)
        if ext in ("xlsx", "xlsm"):
            import io, openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True, read_only=True)
            out = []
            for ws in wb.worksheets:
                out.append(f"# Sheet: {ws.title}")
                rows = 0
                for row in ws.iter_rows(values_only=True):
                    out.append(" | ".join("" if v is None else str(v) for v in row))
                    rows += 1
                    if rows >= 2000:
                        break
            return "\n".join(out)
        if ext in ("csv", "txt", "json", "md"):
            return _safe_decode_text(data)
        if ext in ("png", "jpg", "jpeg", "tiff", "bmp", "webp"):
            if not ATTACH_OCR:
                return ""
            try:
                from PIL import Image
                import io, pytesseract
                img = Image.open(io.BytesIO(data))
                return pytesseract.image_to_string(img, lang=OCR_LANG)
            except Exception:
                return ""
        return ""
    except Exception:
        return ""

def _load_signature() -> str:
    # Prefer file if provided; else env
    if SIGNATURE_FILE:
        p = Path(SIGNATURE_FILE)
        if p.exists():
            try:
                return p.read_text(encoding="utf-8")
            except Exception:
                pass
    # Fallback to local repository file 'signature.txt' if present
    try:
        local_sig = (Path(__file__).parent / "signature.txt").resolve()
        if local_sig.exists():
            return local_sig.read_text(encoding="utf-8")
    except Exception:
        pass
    return SIGNATURE_ENV or ""

SIGNATURE_RAW = _load_signature().strip()

def _append_signature_html(body_text: str) -> str:
    """
    Build final HTML body:
      <div>body…</div>
      <hr>
      <div>signature…</div>
    """
    if not body_text:
        body_text = ""
    body_text = _strip_subject_lines(body_text)
    body_text = _strip_llm_signature(body_text)

    body_html = _text_to_html(body_text).strip()

    if not SIGNATURE_RAW:
        return f"<div>{body_html}</div>"

    sig_html = _text_to_html(SIGNATURE_RAW.strip())
    sep = '<hr style="border:none;border-top:1px solid #ddd;margin:16px 0;">'
    return f"<div>{body_html}</div>{sep}<div>{sig_html}</div>"

# -------- Microsoft Graph (Application permissions) ----------
# Env required: MS_TENANT_ID, MS_CLIENT_ID, MS_CLIENT_SECRET, MS_USER_ID
import requests
import time
import threading as _threading
from typing import Optional, Tuple

MS_TENANT_ID    = os.getenv("MS_TENANT_ID", "")
MS_CLIENT_ID    = os.getenv("MS_CLIENT_ID", "")
MS_CLIENT_SECRET= os.getenv("MS_CLIENT_SECRET", "")
MS_USER_ID      = os.getenv("MS_USER_ID", "")  # target mailbox (email or GUID)

GRAPH_TOKEN_URL = f"https://login.microsoftonline.com/{MS_TENANT_ID}/oauth2/v2.0/token"
GRAPH_BASE      = "https://graph.microsoft.com/v1.0"

_http = requests.Session()
_http.headers.update({"Accept": "application/json"})

# Token cache with expiration
_token_cache: Optional[Tuple[str, float]] = None  # (token, expires_at)
_token_lock = _threading.Lock()
TOKEN_BUFFER_SECONDS = 300  # Refresh token 5 minutes before expiration

# Retry configuration
MAX_RETRIES = 3
RETRY_DELAY_BASE = 1  # Base delay in seconds for exponential backoff
RATE_LIMIT_DELAY = 60  # Delay when hitting rate limits

def _require_env():
    missing = [k for k, v in {
        "MS_TENANT_ID": MS_TENANT_ID,
        "MS_CLIENT_ID": MS_CLIENT_ID,
        "MS_CLIENT_SECRET": MS_CLIENT_SECRET,
        "MS_USER_ID": MS_USER_ID,
    }.items() if not v]
    if missing:
        raise RuntimeError(f"Missing required env: {', '.join(missing)}")

def _retry_with_backoff(func, *args, **kwargs):
    """Retry a function with exponential backoff and rate limit handling."""
    last_exception = None
    
    for attempt in range(MAX_RETRIES + 1):
        try:
            return func(*args, **kwargs)
        except requests.exceptions.RequestException as e:
            last_exception = e
            status_code = getattr(e.response, 'status_code', None) if hasattr(e, 'response') else None
            
            # Handle rate limiting
            if status_code == 429:
                log(f"Rate limited, waiting {RATE_LIMIT_DELAY}s before retry {attempt + 1}/{MAX_RETRIES}")
                time.sleep(RATE_LIMIT_DELAY)
                continue
            
            # Handle server errors (5xx) and some client errors (4xx)
            if status_code and (status_code >= 500 or status_code in [408, 429]):
                if attempt < MAX_RETRIES:
                    delay = RETRY_DELAY_BASE * (2 ** attempt)  # Exponential backoff
                    log(f"API error {status_code}, retrying in {delay}s (attempt {attempt + 1}/{MAX_RETRIES})")
                    time.sleep(delay)
                    continue
            
            # Don't retry on other errors
            raise e
            
        except Exception as e:
            last_exception = e
            if attempt < MAX_RETRIES:
                delay = RETRY_DELAY_BASE * (2 ** attempt)
                log(f"Unexpected error: {type(e).__name__}: {e}, retrying in {delay}s (attempt {attempt + 1}/{MAX_RETRIES})")
                time.sleep(delay)
                continue
            raise e
    
    # If we get here, all retries failed
    raise last_exception

def get_token() -> str:
    """Get a valid access token, using cache if available and not expired."""
    global _token_cache
    
    # Check if we have a valid cached token
    with _token_lock:
        if _token_cache is not None:
            token, expires_at = _token_cache
            current_time = time.time()
            if current_time < (expires_at - TOKEN_BUFFER_SECONDS):
                return token
            else:
                log(f"Token expires in {expires_at - current_time:.0f}s, refreshing...")
    
    # Get new token
    _require_env()
    data = {
        "client_id": MS_CLIENT_ID,
        "client_secret": MS_CLIENT_SECRET,
        "scope": "https://graph.microsoft.com/.default",
        "grant_type": "client_credentials",
    }
    
    try:
        r = _http.post(GRAPH_TOKEN_URL, data=data, timeout=20)
        if r.status_code >= 400:
            raise RuntimeError(f"token: {r.status_code} {r.text}")
        
        response_data = r.json()
        token = response_data.get("access_token")
        if not token:
            raise RuntimeError("token: missing access_token")
        
        # Cache the token with expiration time
        expires_in = response_data.get("expires_in", 3600)  # Default to 1 hour
        expires_at = time.time() + expires_in
        with _token_lock:
            _token_cache = (token, expires_at)
        
        log(f"New token obtained, expires in {expires_in}s")
        return token
        
    except requests.exceptions.RequestException as e:
        raise RuntimeError(f"token request failed: {type(e).__name__}: {e}")
    except Exception as e:
        raise RuntimeError(f"token error: {type(e).__name__}: {e}")

def _auth(token: str) -> dict:
    return {"Authorization": f"Bearer {token}", "Content-Type": "application/json"}

def list_unread_messages(token: str) -> List[Dict[str, Any]]:
    def _make_request():
        url = f"{GRAPH_BASE}/users/{MS_USER_ID}/mailFolders/Inbox/messages"
        params = {
            "$filter": "isRead eq false",
            "$orderby": "receivedDateTime desc",
            "$top": "10",
            "$select": "id,subject,bodyPreview,receivedDateTime,from,body,hasAttachments",
        }
        r = _http.get(url, headers=_auth(token), params=params, timeout=30)
        if r.status_code >= 400:
            raise RuntimeError(f"list_unread: {r.status_code} {r.text}")
        return r.json().get("value", [])
    
    return _retry_with_backoff(_make_request)

def list_attachments(msg_id: str, token: str) -> List[Dict[str, Any]]:
    def _make_request():
        url = f"{GRAPH_BASE}/users/{MS_USER_ID}/messages/{msg_id}/attachments"
        # Let Graph return default fields (including @odata.type); just cap results
        params = {"$top": str(ATTACH_MAX_COUNT)}
        r = _http.get(url, headers=_auth(token), params=params, timeout=30)
        if r.status_code >= 400:
            raise RuntimeError(f"attachments: {r.status_code} {r.text}")
        items = r.json().get("value", [])
        # Filter out reference/item attachments; accept only fileAttachment
        out = []
        for a in items:
            otype = a.get("@odata.type", "")
            if ("fileAttachment" in otype) or (otype == ""):
                out.append(a)
        log(f"[attach] found {len(out)} file attachments out of {len(items)} total")
        return out
    return _retry_with_backoff(_make_request)

def fetch_attachment_bytes(msg_id: str, att_id: str, token: str) -> Dict[str, Any]:
    def _make_request():
        url = f"{GRAPH_BASE}/users/{MS_USER_ID}/messages/{msg_id}/attachments/{att_id}"
        params = {"$select": "name,contentType,size,contentBytes"}
        r = _http.get(url, headers=_auth(token), params=params, timeout=30)
        if r.status_code >= 400:
            raise RuntimeError(f"attach_get: {r.status_code} {r.text}")
        j = r.json() or {}
        b64 = j.get("contentBytes") or ""
        try:
            data = base64.b64decode(b64) if b64 else b""
        except Exception:
            data = b""
        size = int(j.get("size") or 0)
        # Fallback: if no inline contentBytes, download raw via $value (cap by ATTACH_MAX_MB)
        if (not data) and size > 0:
            val_url = f"{GRAPH_BASE}/users/{MS_USER_ID}/messages/{msg_id}/attachments/{att_id}/$value"
            rv = _http.get(val_url, headers=_auth(token), timeout=60, stream=True)
            if rv.status_code >= 400:
                raise RuntimeError(f"attach_value: {rv.status_code} {rv.text}")
            limit = ATTACH_MAX_MB * 1024 * 1024
            buf = bytearray()
            read = 0
            for chunk in rv.iter_content(65536):
                if not chunk:
                    break
                read += len(chunk)
                if read > limit:
                    # stop reading beyond limit
                    break
                buf.extend(chunk)
            data = bytes(buf)
        return {"name": j.get("name",""), "contentType": j.get("contentType",""), "size": size, "bytes": data}
    return _retry_with_backoff(_make_request)

def fetch_message(msg_id: str, token: str) -> Dict[str, Any]:
    def _make_request():
        url = f"{GRAPH_BASE}/users/{MS_USER_ID}/messages/{msg_id}"
        params = {"$select": "id,subject,bodyPreview,receivedDateTime,from,body,conversationId"}
        r = _http.get(url, headers=_auth(token), params=params, timeout=30)
        if r.status_code >= 400:
            raise RuntimeError(f"fetch_message: {r.status_code} {r.text}")
        return r.json()
    
    return _retry_with_backoff(_make_request)

def list_messages_in_conversation(conversation_id: str, token: str, top: int = 10) -> List[Dict[str, Any]]:
    def _make_request():
        url = f"{GRAPH_BASE}/users/{MS_USER_ID}/messages"
        params = {
            "$filter": f"conversationId eq '{conversation_id}'",
            "$orderby": "receivedDateTime desc",
            "$top": str(int(top)),
            "$select": "id,receivedDateTime,hasAttachments"
        }
        r = _http.get(url, headers=_auth(token), params=params, timeout=30)
        if r.status_code >= 400:
            raise RuntimeError(f"list_conv: {r.status_code} {r.text}")
        return r.json().get("value", [])
    return _retry_with_backoff(_make_request)

def mark_read(msg_id: str, token: str) -> None:
    def _make_request():
        url = f"{GRAPH_BASE}/users/{MS_USER_ID}/messages/{msg_id}"
        payload = {"isRead": True}
        r = _http.patch(url, headers=_auth(token), json=payload, timeout=30)
        if r.status_code >= 400:
            raise RuntimeError(f"mark_read: {r.status_code} {r.text}")
    
    _retry_with_backoff(_make_request)

def reply_in_thread(msg_id: str, html_body: str, token: str, reply_all: bool=True) -> None:
    """
    Keep the thread by using the 'reply draft' flow:
      POST   /users/{id}/messages/{msg_id}/createReply[All]
      PATCH  /users/{id}/messages/{draft_id}   (body.contentType='HTML', body.content=html)
      POST   /users/{id}/messages/{draft_id}/send
    """
    def _create_reply():
        base = f"{GRAPH_BASE}/users/{MS_USER_ID}/messages/{msg_id}"
        create_url = base + ("/createReplyAll" if reply_all else "/createReply")
        r = _http.post(create_url, headers=_auth(token), timeout=30)
        if r.status_code >= 400:
            raise RuntimeError(f"reply_create: {r.status_code} {r.text}")
        return r.json() or {}
    
    def _patch_draft(draft_id: str):
        patch_url = f"{GRAPH_BASE}/users/{MS_USER_ID}/messages/{draft_id}"
        # Fetch original message body to ensure thread is visible below our reply
        try:
            orig = fetch_message(msg_id, token) or {}
            orig_html = (((orig.get("body") or {}).get("content") or ""))
        except Exception:
            orig_html = ""
        combined_html = f"<div>{html_body or ''}</div>" + (f"<br>{orig_html}" if orig_html else "")
        payload = {"body": {"contentType": "HTML", "content": combined_html}}
        r = _http.patch(patch_url, headers=_auth(token), json=payload, timeout=30)
        if r.status_code >= 400:
            raise RuntimeError(f"reply_patch: {r.status_code} {r.text}")
    
    def _send_draft(draft_id: str):
        send_url = f"{GRAPH_BASE}/users/{MS_USER_ID}/messages/{draft_id}/send"
        r = _http.post(send_url, headers=_auth(token), timeout=30)
        if r.status_code not in (202, 204):
            raise RuntimeError(f"reply_send: {r.status_code} {r.text}")
    
    # Execute the three-step process with retry logic
    draft = _retry_with_backoff(_create_reply)
    draft_id = draft.get("id")
    if not draft_id:
        raise RuntimeError("reply_create: missing draft id in response")
    
    _retry_with_backoff(_patch_draft, draft_id)
    _retry_with_backoff(_send_draft, draft_id)

# -------- State (de-dupe) ----------
class State:
    def __init__(self, path: Path):
        self.path = path
        self.seen: set[str] = set()
        self.total_replied = 0
        self.last_error: Optional[str] = None
        self._load()

    def _load(self):
        if self.path.exists():
            try:
                data = json.loads(self.path.read_text())
                self.seen = set(data.get("seen", []))
                self.total_replied = int(data.get("total_replied", 0))
                self.last_error = data.get("last_error")
            except Exception:
                pass

    def save(self):
        try:
            # Create backup of existing state
            if self.path.exists():
                backup_path = self.path.with_suffix('.json.backup')
                self.path.rename(backup_path)
            
            # Write new state
            state_data = {
                "seen": sorted(self.seen),
                "total_replied": self.total_replied,
                "last_error": self.last_error,
                "last_saved": datetime.now(timezone.utc).isoformat(),
            }
            self.path.write_text(json.dumps(state_data, indent=2))
            
            # Remove backup on success
            backup_path = self.path.with_suffix('.json.backup')
            if backup_path.exists():
                backup_path.unlink()
                
        except Exception as e:
            log(f"[ERROR] Failed to save state: {type(e).__name__}: {e}")
            # Try to restore backup if it exists
            backup_path = self.path.with_suffix('.json.backup')
            if backup_path.exists():
                try:
                    backup_path.rename(self.path)
                    log("Restored state from backup")
                except Exception as restore_e:
                    log(f"[ERROR] Failed to restore backup: {type(restore_e).__name__}: {restore_e}")
            raise

STATE = State(STATE_PATH)

# -------- Prompt builder ----------
def build_prompt_with_memory(message: Dict[str, Any]) -> str:
    # Validate sender
    if not validate_sender(message.get("from", {})):
        raise ValueError("Invalid sender information")
    
    sender = (((message.get("from") or {}).get("emailAddress") or {}).get("address") or "").lower()
    subject = validate_email_content((message.get("subject") or "").strip())
    body_html = (message.get("body") or {}).get("content") or ""
    
    if (message.get("body") or {}).get("contentType","").lower() == "html":
        body = html_to_text(body_html)
    else:
        body = body_html or (message.get("bodyPreview") or "")
    
    # Validate and sanitize body content
    body = validate_email_content(body)

    # Attachments summary (non-TEACH path): include short snippets to help the model
    attach_block = ""
    if ATTACH_ENABLE:
        try:
            token = get_token()
            # Collect attachments from the current message and, if empty, from recent messages in the same conversation
            atts = list_attachments(message.get("id"), token)
            if not atts:
                try:
                    conv_id = (message.get("conversationId") or "")
                except Exception:
                    conv_id = ""
                if conv_id:
                    for mm in list_messages_in_conversation(conv_id, token, top=5):
                        for a in list_attachments(mm.get("id"), token):
                            atts.append(a)
            pieces = []
            total = 0
            for a in atts[:ATTACH_MAX_COUNT]:
                name = a.get("name") or "attachment"
                ctype = a.get("contentType") or ""
                size = int(a.get("size") or 0)
                inline = bool(a.get("isInline"))
                if inline:
                    log(f"[attach] skip inline: {name} size={size} type={ctype}")
                    continue
                det = fetch_attachment_bytes(message.get("id"), a.get("id"), token)
                data = det.get("bytes", b"")
                if not data:
                    log(f"[attach] empty bytes after fetch (maybe too large or restricted): {name} size={size} type={ctype}")
                snippet = _extract_text_from_attachment(det.get("name",name), det.get("contentType",ctype), data)
                if not snippet: continue
                snippet = snippet.strip().replace("\r\n","\n")
                if len(snippet) > 400: snippet = snippet[:400] + "…"
                entry = f"- {name}: {snippet}"
                pieces.append(entry)
                total += len(entry)
                if total >= ATTACH_SUMMARY_MAX_CHARS: break
            if pieces:
                attach_block = "Attachments:\n" + "\n".join(pieces) + "\n\n"
            else:
                log("[attach] no usable attachments or extraction produced no text")
        except Exception as e:
            log(f"[WARN] attachment summary failed: {type(e).__name__}: {e}")

    # semantic pulls with light filters
    snippets = []
    try:
        snippets = search_memory(query=f"{subject}\n{body[:500]}", k=6, where=None)
    except Exception as e:
        STATE.last_error = f"memory search failed: {e}"

    memory_block = ""
    if snippets:
        lines = []
        for s in snippets[:6]:
            meta = s.get("meta") or {}
            tag_str = meta.get("tags") or ""
            kind = meta.get("kind") or "note"
            lines.append(f"- ({kind}; tags={tag_str}) {s.get('text','')}")
        memory_block = "Relevant internal notes:\n" + "\n".join(lines) + "\n\n"

    # IMPORTANT: tell the model not to generate its own Subject or signature.
    prompt = f"""You are an email assistant. Write a helpful, concise reply to the email below.

Guidelines:
- Write only the email body. Do NOT include a Subject line.
- Do NOT include a full signature block; I'll append mine automatically.
- Keep it professional and brief unless the sender explicitly asks for detail.

{memory_block}{attach_block}Incoming email (from {sender}) — SUBJECT: {subject}

Email body:
{body}
"""
    if len(prompt) > MAX_PROMPT_CHARS:
        prompt = prompt[:MAX_PROMPT_CHARS] + "\n\n[truncated]"
    return prompt

# -------- [TEACH] path ----------
def maybe_handle_teach(m, token) -> bool:
    subj = (m.get("subject") or "").strip()
    if not subj.lower().startswith("[teach]"):
        return False

    sender = (((m.get("from") or {}).get("emailAddress") or {}).get("address") or "").lower()
    body = (m.get("body") or {}).get("content") or m.get("bodyPreview") or ""
    if (m.get("body") or {}).get("contentType","").lower() == "html":
        body = html_to_text(body)

    kind = "rule"
    tags = []
    parsed_text = None
    for line in body.splitlines():
        L = line.strip()
        up = L.upper()
        if up.startswith("KIND:"):
            kind = L.split(":",1)[1].strip().lower() or "rule"
        elif up.startswith("TAGS:"):
            tags = [t.strip() for t in L.split(":",1)[1].split(",") if t.strip()]
        elif up.startswith("TEXT:"):
            parsed_text = L.split(":",1)[1].strip()

    text_to_store = parsed_text or body.strip()
    if not text_to_store and not (ATTACH_ENABLE and (m.get("hasAttachments") or False)):
        reply_in_thread(m["id"], _append_signature_html("I didn't find any content to learn."), token)
        mark_read(m["id"], token)
        return True

    saved_any = False
    # Save body text if present
    if text_to_store:
        try:
            add_memory(text=text_to_store, kind=kind, tags=tags, author=sender, source="email")
            saved_any = True
        except Exception as e:
            log(f"[WARN] failed to save TEACH body: {type(e).__name__}: {e}")

    # Save attachments if enabled
    if ATTACH_ENABLE:
        try:
            atts = list_attachments(m["id"], token)
            if not atts:
                # Scan recent messages in the same conversation for attachments
                try:
                    conv_id = (m.get("conversationId") or "")
                except Exception:
                    conv_id = ""
                if conv_id:
                    for mm in list_messages_in_conversation(conv_id, token, top=5):
                        for a in list_attachments(mm.get("id"), token):
                            atts.append(a)
            count = 0
            for a in atts[:ATTACH_MAX_COUNT]:
                name = a.get("name") or "attachment"
                ctype = a.get("contentType") or ""
                ext = _file_ext(name, ctype)
                if ATTACH_EXTS and ext not in ATTACH_EXTS:
                    log(f"[attach] skip ext not allowed: {name} ext={ext}")
                    continue
                det = fetch_attachment_bytes(m["id"], a.get("id"), token)
                txt = _extract_text_from_attachment(det.get("name", name), det.get("contentType", ctype), det.get("bytes", b""))
                if not txt:
                    log(f"[attach] no text extracted: {name} type={ctype}")
                    continue
                try:
                    add_memory(text=txt, kind=kind, tags=(tags + ["attachment", f"file:{name}", f"ext:{ext}"]), author=sender, source="email-attachment")
                    saved_any = True
                    count += 1
                except Exception as e:
                    log(f"[WARN] failed to save attachment memory: {type(e).__name__}: {e}")
        except Exception as e:
            log(f"[WARN] attachment ingestion failed: {type(e).__name__}: {e}")

    if saved_any:
        reply_in_thread(m["id"], _append_signature_html("Thanks — I’ve saved this to my memory."), token)
    else:
        reply_in_thread(m["id"], _append_signature_html("Sorry — I couldn’t extract any content to save."), token)
    mark_read(m["id"], token)
    return True

# -------- Logging & signals ----------
def log(*args, **kwargs):
    ts = datetime.now(timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ")
    print(ts, *args, flush=True, **kwargs)

_STOP = False
def _sig_handler(signum, frame):
    global _STOP
    _STOP = True

# -------- Health server ----------
_httpd = None
def start_health_server():
    import http.server
    from http.server import ThreadingHTTPServer

    class Handler(http.server.BaseHTTPRequestHandler):
        def _write(self, code: int, ctype: str, payload: bytes):
            self.send_response(code)
            self.send_header("Content-Type", ctype)
            self.send_header("Content-Length", str(len(payload)))
            self.end_headers()
            self.wfile.write(payload)

        def do_GET(self):
            if self.path == "/health":
                # Check external dependencies
                health_status = "ok"
                issues = []
                
                # Check if we can get a token
                try:
                    test_token = get_token()
                    if not test_token:
                        health_status = "degraded"
                        issues.append("token_retrieval_failed")
                except Exception as e:
                    health_status = "degraded"
                    issues.append(f"token_error: {type(e).__name__}")
                
                # Check memory system availability without embeddings
                try:
                    if _HAS_MEM and callable(mem_all_stats):
                        _ = mem_all_stats()
                    elif not _HAS_MEM:
                        health_status = "degraded"
                        issues.append("memory_system_unavailable")
                except Exception as e:
                    health_status = "degraded"
                    issues.append(f"memory_error: {type(e).__name__}")
                
                body = json.dumps({
                    "status": health_status,
                    "now": datetime.now(timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ"),
                    "dry_run": DRY_RUN,
                    "total_replied": STATE.total_replied,
                    "processed_cache": len(STATE.seen),
                    "last_error": STATE.last_error,
                    "poll_seconds": POLL_SECONDS,
                    "issues": issues,
                    "uptime_seconds": int(time.time() - _start_time) if '_start_time' in globals() else 0,
                }).encode("utf-8")
                
                status_code = 200 if health_status == "ok" else 503
                self._write(status_code, "application/json; charset=utf-8", body)
            elif self.path == "/memstats":
                if _HAS_MEM and callable(mem_all_stats):
                    try:
                        payload = json.dumps(mem_all_stats()).encode("utf-8")
                        self._write(200, "application/json; charset=utf-8", payload)
                    except Exception as e:
                        msg = json.dumps({"error": f"memstats failed: {type(e).__name__}: {e}"}).encode("utf-8")
                        self._write(500, "application/json; charset=utf-8", msg)
                else:
                    msg = json.dumps({"error": "memory module not available"}).encode("utf-8")
                    self._write(501, "application/json; charset=utf-8", msg)
            elif self.path.startswith("/memdump"):
                try:
                    from urllib.parse import urlparse, parse_qs
                    if not _HAS_MEM:
                        msg = json.dumps({"error": "memory module not available"}).encode("utf-8")
                        self._write(501, "application/json; charset=utf-8", msg)
                        return
                    u = urlparse(self.path)
                    qs = parse_qs(u.query or "")
                    q = (qs.get("q", [""])[0] or "").strip()
                    # Parse limit robustly (supports values like "50.") and clamp to sane range
                    _raw_limit = (qs.get("limit", ["50"])[0] or "50").strip()
                    try:
                        limit = int(float(_raw_limit))
                    except Exception:
                        limit = 50
                    limit = max(1, min(limit, 500))
                    kind = (qs.get("kind", [None])[0] or None)
                    tag_contains = (qs.get("tag_contains", [None])[0] or None)

                    where = {}
                    if kind:
                        where["kind"] = kind
                    if tag_contains:
                        where["tags"] = {"$contains": tag_contains}

                    items = []
                    if q:
                        items = search_memory(query=q, k=limit, where=where)
                    else:
                        try:
                            from core.vector_mem import _collection
                            # 'ids' is always returned; include controls optional fields
                            res = _collection.get(limit=limit, offset=0, include=["documents","metadatas"])
                            ids = res.get("ids", []) or []
                            docs = res.get("documents", []) or []
                            metas = res.get("metadatas", []) or []
                            for i in range(len(ids)):
                                items.append({"id": ids[i], "text": (docs[i] if i < len(docs) else ""), "meta": (metas[i] if i < len(metas) else {})})
                        except Exception as e:
                            msg = json.dumps({"error": f"raw get failed: {type(e).__name__}: {e}"}).encode("utf-8")
                            self._write(500, "application/json; charset=utf-8", msg)
                            return

                    body = json.dumps({"count": len(items), "items": items}).encode("utf-8")
                    self._write(200, "application/json; charset=utf-8", body)
                except Exception as e:
                    msg = json.dumps({"error": f"memdump failed: {type(e).__name__}: {e}"}).encode("utf-8")
                    self._write(500, "application/json; charset=utf-8", msg)
            elif self.path in ("/", "/index.html"):
                html_doc = f"""<!doctype html>
<html>
  <head><meta charset="utf-8"><title>Simple Email Assistant</title></head>
  <body style="font-family: system-ui, -apple-system, Segoe UI, Roboto, sans-serif; margin: 2rem;">
    <h1>Simple Email Assistant</h1>
    <p>Status: <strong>running</strong></p>
    <ul>
      <li>Dry run: <code>{DRY_RUN}</code></li>
      <li>Total replied: <code>{STATE.total_replied}</code></li>
      <li>Processed cache: <code>{len(STATE.seen)}</code></li>
      <li>Health: <a href="/health">/health</a></li>
      <li>Mem stats: <a href="/memstats">/memstats</a></li>
      <li>Mem dump: <a href="/memdump">/memdump</a></li>
    </ul>
  </body>
</html>"""
                self._write(200, "text/html; charset=utf-8", html_doc.encode("utf-8"))
            else:
                self._write(404, "text/plain; charset=utf-8", b"Not Found\n")
        def log_message(self, fmt, *args): pass

    class QuietTCPServer(ThreadingHTTPServer):
        daemon_threads = True
        allow_reuse_address = True

    global _httpd
    _httpd = QuietTCPServer(("0.0.0.0", HEALTH_PORT), Handler)
    t = threading.Thread(target=_httpd.serve_forever, daemon=True)
    t.start()
    log(f"[health] listening on :{HEALTH_PORT}")

def stop_health_server():
    global _httpd
    if _httpd is not None:
        try:
            _httpd.shutdown()
            _httpd.server_close()
        except Exception:
            pass
        _httpd = None

# -------- Main worker loop ----------
def _sender_allowed(m: Dict[str, Any]) -> bool:
    if not (ALLOWLIST_EMAILS or ALLOWLIST_DOMAINS):
        return True
    sender = (((m.get("from") or {}).get("emailAddress") or {}).get("address") or "").lower()
    if sender in ALLOWLIST_EMAILS:
        return True
    dom = sender.split("@")[-1] if "@" in sender else ""
    return dom in ALLOWLIST_DOMAINS

def _is_self_email(m: Dict[str, Any]) -> bool:
    """Check if the email is from the assistant itself to prevent loops."""
    sender = (((m.get("from") or {}).get("emailAddress") or {}).get("address") or "").lower()
    # Check if sender matches the target mailbox
    return sender == MS_USER_ID.lower() or sender == f"noreply@{MS_USER_ID.split('@')[-1] if '@' in MS_USER_ID else ''}"

def run_worker_loop():
    import traceback
    token = None
    while not _STOP:
        try:
            token = get_token()
            msgs = list_unread_messages(token)
            if msgs:
                log(f"Found {len(msgs)} unread message(s).")

            for m in msgs or []:
                # De-dupe: skip if already processed
                mid = m.get("id")
                if mid and mid in STATE.seen:
                    continue
                if not _sender_allowed(m):
                    # leave unread; skip completely
                    continue
                
                # Prevent email loops
                if _is_self_email(m):
                    log(f"Skipping self-email from {(((m.get('from') or {}).get('emailAddress') or {}).get('address') or '?')}")
                    try:
                        mark_read(m["id"], token)
                    except Exception as mark_e:
                        log(f"[ERROR] Failed to mark self-email as read: {type(mark_e).__name__}: {mark_e}")
                    continue

                # Try TEACH flow first
                try:
                    if maybe_handle_teach(m, token):
                        continue
                except Exception as e:
                    log(f"[ERROR] teach handler failed: {type(e).__name__}: {e}")
                    try: 
                        mark_read(m["id"], token)
                    except Exception as mark_e:
                        log(f"[ERROR] Failed to mark message as read: {type(mark_e).__name__}: {mark_e}")
                    continue

                # Normal answer flow
                try:
                    prompt = build_prompt_with_memory(m)
                    ans_text = ask_chatgpt(prompt) or ""
                    final_html = _append_signature_html(ans_text)

                    if not DRY_RUN:
                        reply_in_thread(m["id"], final_html, token, reply_all=True)
                        mark_read(m["id"], token)

                    STATE.total_replied += 1
                    STATE.seen.add(m["id"])
                    STATE.save()

                    sender_addr = (((m.get("from") or {}).get("emailAddress") or {}).get("address") or "?")
                    log(f"Replied to {sender_addr} msg_id={m.get('id')}")
                except Exception as e:
                    STATE.last_error = f"{type(e).__name__}: {e}"
                    STATE.save()
                    log(f"[ERROR] {STATE.last_error}\n{traceback.format_exc()}")

        except Exception as e:
            import traceback as _tb
            STATE.last_error = f"top loop: {type(e).__name__}: {e}"
            STATE.save()
            log(f"[ERROR] {STATE.last_error}\n{_tb.format_exc()}")

        # Sleep with fast shutdown responsiveness
        for _ in range(POLL_SECONDS):
            if _STOP: break
            time.sleep(1)

if __name__ == "__main__":
    signal.signal(signal.SIGINT, _sig_handler)
    signal.signal(signal.SIGTERM, _sig_handler)

    # Track startup time for health monitoring
    global _start_time
    _start_time = time.time()

    print("[boot] starting health server...", flush=True)
    start_health_server()
    print("[boot] health server started", flush=True)

    print(f"{datetime.now(timezone.utc).strftime('%Y-%m-%dT%H:%M:%SZ')} Email assistant started. poll={POLL_SECONDS}s dry_run={int(DRY_RUN)}", flush=True)
    run_worker_loop()
    print("Shutting down. Bye.", flush=True)
